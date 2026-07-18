import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # We want 16:9 ratio. Default is 4:3.
    # 16:9 in inches is typically 13.33 x 7.5
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Layouts
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # Helper function to format title
    def format_title(slide, text):
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.name = 'Arial'
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(27, 58, 92) # Dark blue

    # Helper function to add bullets
    def add_bullets(slide, items):
        body = slide.placeholders[1]
        tf = body.text_frame
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.font.name = 'Arial'

    # Slide 1: Title
    slide = prs.slides.add_slide(title_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "NTTH: No Time To Hack"
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(27, 58, 92)
    
    subtitle.text = "An Agent-Inspired Autonomous Network Defense Architecture\nwith Hybrid Risk Scoring and Dynamic Honeypot Deployment"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(85, 85, 85)

    # Slide 2: Table of Contents
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Table of Contents")
    add_bullets(slide, [
        "1. Introduction & Problem Statement",
        "2. What is NTTH?",
        "3. System Architecture",
        "4. How it Works (Core Principles)",
        "5. AI Agents Workflow",
        "6. Threat & Decision Stages",
        "7. Enforcement & Reporting Stages",
        "8. Working of AR9271 WiFi Monitor",
        "9. Current Working & Deployed Stage",
        "10. Expected Outcomes",
        "11. Future Work",
        "12. Conclusion"
    ])

    # Slide 3: Introduction & Problem Statement
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Introduction & Problem Statement")
    add_bullets(slide, [
        "Automated attacks act at machine speed.",
        "Traditional SOCs involve human analysts who take 15 to 30 minutes to investigate and respond.",
        "This manual process creates a critical gap between detection and containment.",
        "Traditional IDS (Snort, Suricata) only generate alerts and lack automated response capabilities.",
        "Goal: Eliminate the human-in-the-loop delay and close the 15-30 minute window."
    ])

    # Slide 4: What is NTTH?
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "What is NTTH?")
    add_bullets(slide, [
        "NTTH (No Time To Hack) is an autonomous network defense system.",
        "Implements a complete closed-loop pipeline: Capture -> Detect -> Decide -> Enforce -> Trap.",
        "Achieves sub-second response latency (~127 ms).",
        "Utilizes a hybrid detection model (Rule-based + ML Anomaly Detection).",
        "Dynamically deploys Linux kernel firewall rules (nftables).",
        "Transparently redirects attackers to deception surfaces (Cowrie SSH, HTTP Honeypot)."
    ])

    # Slide 5: System Architecture
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "System Architecture")
    add_bullets(slide, [
        "A five-phase modular pipeline connected via an asynchronous event bus (publish/subscribe).",
        "Phase 1: Network Capture (Scapy Sniffer + AR9271 Monitor).",
        "Phase 2: Event Bus for decoupled inter-agent communication.",
        "Phase 3: Agent Pipeline (Threat, Decision, Enforcement, Reporting).",
        "Phase 4: Deception Layer (Honeypot redirection).",
        "Phase 5: Real-Time Cross-Platform Dashboard (Flutter Web + Android)."
    ])

    # Slide 6: How it Works (Core Principles)
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "How it Works (Core Principles)")
    add_bullets(slide, [
        "Packets are continuously intercepted without halting legitimate traffic.",
        "Features are extracted (10-dimensional vector) for each packet.",
        "Traffic is assigned a risk score between 0.0 and 1.0.",
        "Actions are dynamically selected based on risk thresholds.",
        "Threats >= 0.40 trigger dynamic honeypot redirection.",
        "Threats >= 0.80 trigger total firewall blocking."
    ])

    # Slide 7: AI Agents Workflow
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "AI Agents Workflow")
    add_bullets(slide, [
        "Inspired by Russell and Norvig's perception-reasoning-action model.",
        "Agents act as autonomous modules:",
        "1. Threat Agent: Perceives packets, reasons with IDS+ML, publishes risk.",
        "2. Decision Agent: Perceives risk, applies policy rules, selects action.",
        "3. Enforcement Agent: Perceives directives, applies nftables kernel rules.",
        "4. Reporting Agent: Perceives complete event, streams via WebSocket to Dashboard."
    ])

    # Slide 8: Threat & Decision Stages
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Threat & Decision Stages")
    add_bullets(slide, [
        "Threat Stage:",
        "  - Hybrid Model = 0.6 * Rule Score + 0.4 * ML Score.",
        "  - Rules detect known attacks (Port Scan, SYN Flood, Brute Force).",
        "  - Isolation Forest ML detects novel anomalies unsupervised.",
        "Decision Stage:",
        "  - Protocol-aware routing avoids blocking critical gateways.",
        "  - Deduplicates repetitive threat events.",
        "  - Directs traffic: Allow, Log, Rate Limit, Honeypot, or Block."
    ])

    # Slide 9: Enforcement & Reporting Stages
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Enforcement & Reporting Stages")
    add_bullets(slide, [
        "Enforcement Stage:",
        "  - Executes direct commands to Linux kernel via nftables.",
        "  - Flow-Aware Redirects: Creates specific DNAT rules for Attacker -> Victim flows.",
        "  - Leaves legitimate user traffic completely unaffected.",
        "Reporting Stage:",
        "  - Persists data to PostgreSQL database.",
        "  - Enriches threat data with GeoIP locations.",
        "  - Pushes live telemetry to the Flutter UI over persistent WebSockets."
    ])

    # Slide 10: Working of AR9271 WiFi Monitor
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Working of AR9271 WiFi Monitor")
    add_bullets(slide, [
        "Commodity hardware adapter (~$10) operating in Monitor Mode.",
        "Sniffs 802.11 management frames natively from the airwaves.",
        "Extracts Probe Requests to track mobile devices seeking known networks.",
        "Detects Deauthentication attacks in real-time.",
        "Provides physical proximity intelligence without requiring devices to connect.",
        "Complements the wired Scapy sniffer for full spectrum awareness."
    ])

    # Slide 11: Current Working & Deployed Stage
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Current Working & Deployed Stage")
    add_bullets(slide, [
        "System is fully functional, containerized via Docker Compose.",
        "Python backend, Scapy, Postgres, Cowrie, and Flutter Dashboard running smoothly.",
        "End-to-End verified against live tests (nmap, Hydra brute force, hping3 SYN flood).",
        "Dashboard handles WebSocket streams effectively, plotting threats and map topologies.",
        "Currently capturing credentials and attacker commands via the honeypot layer."
    ])

    # Slide 12: Expected Outcomes
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Expected Outcomes")
    add_bullets(slide, [
        "Detection Rate: > 95% True Positive Rate for common attack vectors.",
        "False Positives: < 5% due to hybrid scoring balance.",
        "Response Latency: Maintained consistently under 200 ms.",
        "A demonstrable operational advantage over Snort/Suricata manual response paradigms.",
        "Secure containment of threats while gathering rich adversary intelligence."
    ])

    # Slide 13: Future Work
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Future Work")
    add_bullets(slide, [
        "Feedback Agent: Implement an adaptive learning mechanism based on user feedback to adjust risk thresholds.",
        "Reinforcement Learning: Use RL for advanced dynamic response strategy selection.",
        "LLM Honeypot Integration: Connect LLMs to dynamically generate honeypot responses to keep attackers engaged longer.",
        "Distributed Deployment: Extend architecture across multiple sensors across an enterprise WAN."
    ])

    # Slide 14: Conclusion
    slide = prs.slides.add_slide(bullet_layout)
    format_title(slide, "Conclusion")
    add_bullets(slide, [
        "NTTH successfully transforms network defense from passive alerting to active, autonomous enforcement.",
        "By merging ML anomaly detection with flow-aware deception, attackers are contained transparently.",
        "The agent-inspired asynchronous architecture guarantees high-speed and fault-isolated performance.",
        "Provides enterprise-grade defense mechanisms accessible to small-scale networks using commodity hardware."
    ])
    
    # Final 'Thank You' Slide
    slide = prs.slides.add_slide(title_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(27, 58, 92)
    
    subtitle.text = "Questions?"
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)

    prs.save("/home/ubuntu/NTTH/docs/NTTH_Presentation.pptx")
    print("Presentation saved to /home/ubuntu/NTTH/docs/NTTH_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
